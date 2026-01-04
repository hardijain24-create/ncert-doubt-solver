
import streamlit as st
import os
import PyPDF2
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from io import BytesIO

# =========================
# CONFIG
# =========================
BASE_PDF_DIR = "backend_pdfs"
CHUNK_SIZE = 160
TOP_K = 4

st.set_page_config(page_title="NCERT Doubt Solver", layout="centered")

# =========================
# TEXT HELPERS
# =========================

def clean_text(text):
    return " ".join(text.replace("\n", " ").split())

def chunk_text(text, page, source):
    words = text.split()
    chunks = []
    for i in range(0, len(words), CHUNK_SIZE):
        chunk = " ".join(words[i:i + CHUNK_SIZE])
        if len(chunk) > 100:
            chunks.append({
                "text": chunk,
                "page": page,
                "source": source
            })
    return chunks

# =========================
# BUILD INDEX
# =========================

def build_ncert_index(grade, language, extra_chunks=None):
    folder = os.path.join(BASE_PDF_DIR, f"Grade{grade}", language)
    chunks = []

    if os.path.exists(folder):
        for file in os.listdir(folder):
            if file.lower().endswith(".pdf"):
                path = os.path.join(folder, file)
                try:
                    reader = PyPDF2.PdfReader(open(path, "rb"))
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text()
                        if text:
                            chunks.extend(
                                chunk_text(clean_text(text), i + 1, file)
                            )
                except:
                    pass

    if extra_chunks:
        chunks.extend(extra_chunks)

    texts = [c["text"] for c in chunks]
    vectorizer = TfidfVectorizer(stop_words="english")
    tfidf_matrix = vectorizer.fit_transform(texts) if texts else None

    return chunks, vectorizer, tfidf_matrix

# =========================
# RETRIEVAL
# =========================

def retrieve_chunks(question, chunks, vectorizer, tfidf_matrix):
    if tfidf_matrix is None or vectorizer is None or not chunks:
        return []
    q_vec = vectorizer.transform([question])
    scores = cosine_similarity(q_vec, tfidf_matrix)[0]
    top_indices = scores.argsort()[-TOP_K:][::-1]
    return [chunks[i] for i in top_indices]

# =========================
# ANSWER + SLIDES
# =========================

def build_long_answer_and_slides(retrieved, grade):
    combined = " ".join(r["text"] for r in retrieved)
    sentences = [s.strip() for s in combined.split(".") if len(s.strip()) > 30]

    paragraphs = []
    for i in range(0, min(len(sentences), 12), 3):
        paragraphs.append(". ".join(sentences[i:i+3]) + ".")

    long_answer = "\n\n".join(paragraphs) if paragraphs else "No relevant NCERT content found."

    slides = [sentences[i:i+4] for i in range(0, min(len(sentences), 12), 4)]

    citations = {
        f"Page {r['page']} | NCERT Class {grade} Science"
        for r in retrieved
    }

    return long_answer, slides, citations

def generate_ppt(question, slides_content):
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = question
    title_slide.placeholders[1].text = "NCERT Based Explanation"

    for bullets in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Points"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for b in bullets:
            body.add_paragraph().text = b

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# =========================
# SIDEBAR (with emojis)
# =========================

st.sidebar.title("🎓 Select Class & Subject")
grade = st.sidebar.selectbox("Class", [6, 7, 8, 9, 10])
language = st.sidebar.selectbox("Language", ["English", "Hindi"])

uploaded_files = st.sidebar.file_uploader(
    "📎 Optional: Upload your own PDFs",
    type=["pdf"],
    accept_multiple_files=True
)

def parse_uploaded_pdfs(files):
    chunks = []
    for up in files:
        try:
            reader = PyPDF2.PdfReader(up)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    chunks.extend(
                        chunk_text(clean_text(text), i + 1, up.name)
                    )
        except:
            pass
    return chunks

if "uploaded_chunks" not in st.session_state:
    st.session_state.uploaded_chunks = []

if uploaded_files:
    st.session_state.uploaded_chunks = parse_uploaded_pdfs(uploaded_files)
    st.sidebar.success(f"Parsed {len(st.session_state.uploaded_chunks)} uploaded chunks ✅")

# =========================
# ONE-TIME BUILD (with emojis)
# =========================

if "index_ready" not in st.session_state:
    st.session_state.index_ready = False

if not st.session_state.index_ready:
    st.info("📚 NCERT knowledge base not built yet.")

    if st.button("⚙️ Build NCERT Knowledge Base (One-Time)"):
        with st.spinner("Building NCERT knowledge base... Please wait."):
            chunks, vectorizer, tfidf_matrix = build_ncert_index(
                grade,
                language,
                st.session_state.uploaded_chunks
            )

            if not chunks:
                st.error("No PDFs found in backend_pdfs or from uploads.")
                st.stop()

            st.session_state.chunks = chunks
            st.session_state.vectorizer = vectorizer
            st.session_state.tfidf_matrix = tfidf_matrix
            st.session_state.index_ready = True

        st.success("Knowledge base ready ✅")
        st.rerun()

    st.stop()

# =========================
# MAIN UI (with emojis)
# =========================

st.title("📘 NCERT Doubt Solver")
st.caption("Grades 6–10 | Science | Offline | NCERT Only")

question = st.text_input("💬 Ask your question:", placeholder="Example: What are nutrients?")

if question:
    retrieved = retrieve_chunks(
        question,
        st.session_state.chunks,
        st.session_state.vectorizer,
        st.session_state.tfidf_matrix
    )

    if not retrieved:
        st.warning("⚠️ This question is not covered in NCERT.")
    else:
        answer, slides_content, citations = build_long_answer_and_slides(
            retrieved, grade
        )

        st.subheader("📖 Answer")
        st.write(answer)

        st.subheader("📚 Source")
        for c in citations:
            st.write("•", c)

        ppt = generate_ppt(question, slides_content)
        st.download_button(
            "⬇️ Download PPT",
            ppt,
            file_name="NCERT_Answer.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if st.sidebar.button("🔁 Reset Knowledge Base"):
    st.session_state.index_ready = False
    st.session_state.uploaded_chunks = []
    st.rerun()

    ppt = generate_ppt(question, answer, citations)
    with open(ppt, "rb") as f:
        st.download_button("📥 Download PPT", f, file_name="NCERT_Answer.pptx")
